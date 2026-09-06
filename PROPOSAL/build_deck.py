"""Build the Round-2 deck (PPTX) from the slide spec.

Run:  python PROPOSAL/build_deck.py   ->   PROPOSAL/JALE_Round2_deck.pptx

Plain language everywhere; every slide footer carries the source tags.
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
CHARTS = ROOT / "PROPOSAL" / "charts"
OUT = ROOT / "PROPOSAL" / "JALE_Round2_deck.pptx"

NAVY = RGBColor(0x14, 0x34, 0x2B)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
ORANGE = RGBColor(0xE0, 0x7A, 0x5F)
GREY = RGBColor(0x5A, 0x64, 0x6E)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_slide(title: str, tag: str = "") -> object:
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.5), Inches(0.12), Inches(11), Inches(0.7))
    tf = tb.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(26); p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    if tag:
        tg = s.shapes.add_textbox(Inches(10.0), Inches(0.15), Inches(2.9), Inches(0.6))
        tf = tg.text_frame; tf.clear(); p = tf.paragraphs[0]
        p.text = tag; p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(11); p.font.color.rgb = TEAL
    return s


def bullets(s: object, items: list[str], left=0.6, top=1.2, w=11.5, h=4.6,
            size=16):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if item.startswith("__BIG__"):
            p.text = item.replace("__BIG__", "")
            p.font.size = Pt(size + 4); p.font.bold = True; p.font.color.rgb = NAVY
        else:
            p.text = item
            p.font.size = Pt(size)
        p.space_after = Pt(10)
    return tb


def add_image(s: object, path: str, left=7.0, top=1.4, w=5.9):
    s.shapes.add_picture(str(CHARTS / path), Inches(left), Inches(top),
                         width=Inches(w))


def footer(s: object, text: str):
    tb = s.shapes.add_textbox(Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.5))
    tf = tb.text_frame; tf.clear(); p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(10); p.font.color.rgb = GREY


# ---- S1 title ---------------------------------------------------------------
s = add_slide("JA\u00b7LE \u2014 Catch a fraud ring before it is finished forming.",
              "TVS Credit EPIC 8.0 \u00b7 Problem (e)")
bullets(s, [
    "__BIG__Swarm Intelligence Lending Network",
    "",
    "Other teams will tell you a fraud ring exists.",
    "__BIG__JA\u00b7LE tells you ~7 months earlier that one is about to form, "
    "names the one cheap link that breaks it, and proves that blocking "
    "the dealer does nothing.",
    "",
    "Median 204-day lead time \u00b7 Zero false alarms at the strict "
    "operating point \u00b7 Every number regenerable in 30 seconds.",
    "",
    "Pipeline reruns in ~30 s  \u00b7  Synthetic + public data  \u00b7  "
    "Every chart traceable to a committed script",
], size=18)

# ---- S2 macro ---------------------------------------------------------------
s = add_slide("The fraud war TVS fights is a lending war", "[PUBLIC]")
bullets(s, [
    "__BIG__Lending fraud = \u20b940,774 cr = 85% of all reported fraud (FY26)",
    "Card / internet fraud: \u20b929 cr. The real threat is in advances.",
    "And each case is getting bigger: +118% severity per case in 2 years "
    "(\u20b92.17 cr \u2192 \u20b93.83 cr \u2192 \u20b94.72 cr per case).",
    "Source: RBI Annual Report 2025-26.",
], w=6.4)
add_image(s, "g6_rbi.png", left=7.0)
footer(s, "[PUBLIC] RBI AR 2025-26  \u00b7  derived severity/case")

# ---- S3 rings at TVS --------------------------------------------------------
s = add_slide("How a fraud ring actually gets filed at TVS", "[PUBLIC]")
bullets(s, [
    "__BIG__No single application looks suspicious. "
    "The ring lives between records.",
    "Dealer-sourced journey: lead \u2192 KYC \u2192 documents \u2192 "
    "underwriting \u2192 disbursal \u2192 EMI.",
    "",
    "Real ring cases that have hurt NBFCs:",
    "   \u2022 Mahindra Finance: ~\u20b9150 cr, ~2,000 ghost customers, forged KYC",
    "   \u2022 Video-KYC abuse: 35 bikes financed on fabricated papers",
    "   \u2022 Branch-staff collusion with fabricated invoices",
    "",
    "Retailer reality: a DSR (Dealer Sales Rep) enters customer details "
    "into the dealer portal. In rural Tier-3/4 areas with poor "
    "connectivity, video-KYC gaps and off-book guarantor arrangements "
    "create the blind spots fraud rings exploit \u2014 without the "
    "dealer's knowledge.",
    "",
    "By the time a ring is dense enough to flag, the money is "
    "usually already disbursed.",
], size=14)
footer(s, "[PUBLIC] Mahindra Finance (BS/IE/ET Apr-2024) \u00b7 TOI Jan-2024, Apr-2025 \u00b7 retailer journey insight")

# ---- S4 generic trap --------------------------------------------------------
s = add_slide("The generic answer \u2014 and why it fails", "")
bullets(s, [
    "__BIG__Ask any AI tool for a fraud-ring solution and you get "
    "the same three boxes:",
    "Graph + GNN + red-dot dashboard. ~400 teams will show variants "
    "of this.",
    "",
    "Its two failures:",
    "   1. LATE \u2014 by the time a ring is dense enough to detect, "
    "the money is already gone.",
    "   2. ACCUSATORY \u2014 it blames the most-connected node (a busy dealer, "
    "a village shop-phone), punishing genuine rural customers.",
    "",
    "__BIG__We ask a different question:",
    "Which ecosystem is about to form? What is the one cheap link "
    "that causally breaks it? And why does blocking the dealer "
    "change nothing?",
], size=16)
footer(s, "")

# ---- S5 what we built -------------------------------------------------------
s = add_slide("What we built \u2014 five layers, one picture", "")
bullets(s, [
    "__BIG__Data \u2192 Link records \u2192 Score clusters \u2192 "
    "Explain \u2192 Decide",
    "",
    "L1  Entity Resolution \u2014 probabilistic identity linkage "
    "(Fellegi\u2013Sunter, no labels, EM-derived threshold).",
    "L2  Heterogeneous Graph \u2014 5 relations: device / account / "
    "person / guarantor / dealer as sparse incidence matrices.",
    "L3  Application Scoring \u2014 86 features + gradient-boosted trees, "
    "ring-disjoint cross-validation.",
    "L4  Ring-Level Scoring \u2014 score the cluster, not the person. "
    "Fixed structural formula, zero learned weights.",
    "L5  Explanation \u2014 every flag ships an evidence case-file "
    "an analyst can check against the loan file.",
    "",
    "__BIG__Synthetic + public data \u00b7 ~30 s to rerun \u00b7 "
    "every number traceable to a script.",
], size=14)
footer(s, "Demo: jale_demo.html (single self-contained file)")

# ---- S6 receipts ------------------------------------------------------------
s = add_slide("What we measure \u2014 the receipt", "[COLAB]")
bullets(s, [
    "__BIG__Ring-detection AUCPR: 0.754  (20\u00d7 lift over random)",
    "   \u2022 Honest range with nested CV: 0.68 \u2013 0.75",
    "   \u2022 GraphSAGE GNN on the SAME folds & features: 0.688",
    "   \u2022 Best single column alone: 0.233 (6.2\u00d7) \u2014 our honest floor",
    "",
    "The graph features, not the exotic model, carry the signal "
    "\u2014 and we prove it by running every model on the same folds.",
    "",
    "Validated on a real public graph (YelpChi, 45,954 nodes): "
    "leakage-gap protocol transfers, measured at +0.022 AUCPR.",
], w=6.4, size=15)
add_image(s, "g4_models.png", left=7.0)
footer(s, "[COLAB] executed notebook §4/§5/§6b · SMOKE, 3,474 apps, 131 fraud · YelpChi real-graph validation")

# ---- S7 honesty -------------------------------------------------------------
s = add_slide("Why you can trust the number", "[COLAB + SIM]")
bullets(s, [
    "__BIG__Three controls, all passing:",
    "   • Shuffled labels → model drops to chance (no memorisation)",
    "   • Ring-disjoint CV → test rings never leak into training",
    "   • Nested CV → hyperparameters chosen inside inner folds only",
    "",
    "A naive split inflates AUCPR by +10% on synthetic, +2.2% "
    "on a real graph (YelpChi). We report the harder number.",
    "",
    "__BIG__And we tell you where we fail:",
    "Hold out one typology and AUCPR drops to 0.27 — so we "
    "flag novel structures (a built-in novelty detector), "
    "not claim we understand all rings.",
], w=6.4, size=14)
add_image(s, "g2_leakage.png", left=7.0)
footer(s, "[COLAB] leakage gap synthetic +0.100, real +0.022 \u00b7 typology hold-out 0.27 [REPO]")

# ---- S8 differentiator ------------------------------------------------------
s = add_slide("Before fraud occurs \u2014 the measured lead time", "[SIM]")
bullets(s, [
    "__BIG__Trained only on the past, the engine flags 5/9 "
    "forming rings before their last application arrives.",
    "",
    "   \u2022 Strict setting (L4 \u2265 0.70, cluster \u2265 6):",
    "     0 false alarms, median lead 204 days (~7 months)",
    "   \u2022 Loose setting: catches 9/9 at ~5 reviews/week "
    "\u2014 an honest trade-off, not a magic number",
    "",
    "This is literally the problem statement:",
    "__BIG__\"predict emerging fraud ecosystems before fraud occurs.\"",
], w=6.4, size=15)
add_image(s, "g1_lead_time.png", left=7.0)
footer(s, "[SIM] experiments/lead_time.py \u00b7 warm-start on first half of timeline, watch second half")

# ---- S9 which link ----------------------------------------------------------
s = add_slide("Which connection actually carries the risk?", "[COLAB]")
bullets(s, [
    "__BIG__We test every relationship by removing it and measuring "
    "what happens:",
    "",
    "   \u2022 Remove bank-account signal \u2192 detection drops most "
    "(\u22120.22 AUCPR)",
    "   \u2022 Remove device signal \u2192 \u22120.21",
    "   \u2022 Remove dealer signal \u2192 \u22120.10 (LEAST impact)",
    "",
    "__BIG__The dealer is a conduit, not the cause.",
    "   Blocking the dealer changes almost nothing about "
    "the ring's detection. Blocking the bank-account link "
    "collapses it.",
    "",
    "Design consequence: verify the cheap link (account, device), "
    "never blanket-block the dealer. This protects genuine "
    "rural borrowers.",
], w=6.4, size=14)
add_image(s, "g3_ablation.png", left=7.1)
footer(s, "[COLAB] notebook \u00a78a \u00b7 fixed-param node+graph GBT, ring-disjoint")

# ---- S10 decision layer -----------------------------------------------------
s = add_slide("Ring-vaccination: the cheapest link that breaks the ring",
              "[SIM]")
bullets(s, [
    "__BIG__Score the cluster, not the person.",
    "   \u2022 Ring-level queue with an evidence case-file per cluster",
    "   \u2022 Dominant axis shown: device / account / guarantor / burst",
    "",
    "__BIG__Actions are firebreaks \u2014 minimal-friction, "
    "budget-aware:",
    "   \u2022 Step-up: verify one disbursement account",
    "   \u2022 Verify the guarantor pool",
    "   \u2022 Hold disbursement to one account",
    "   \u2022 Dealer audit only when implicated by the ablation",
    "",
    "Every action maps to a measured signal drop: "
    "account \u22120.22 \u00b7 device \u22120.21 \u00b7 "
    "dealer \u22120.10 (least signal).",
    "",
    "Vaccinate the cheapest link. The dealer relationship "
    "and honest borrowers stay untouched.",
], w=6.4, size=14)
add_image(s, "g7_budget.png", left=7.0)
footer(s, "[SIM] repo L4 evaluation: all 12 rings flagged at \u22650.5, 83% fraud apps covered")

# ---- S11 business case ------------------------------------------------------
s = add_slide("The business case \u2014 anchored to TVS\u2019s audited books",
              "[PUBLIC + A#]")
bullets(s, [
    "__BIG__Year-1 enterprise investment \u2192 Year-2 breakeven \u2192 Year-3 profit.",
    "",
    "Audited anchors (AR Note 18, FY26):",
    "   \u2022 Target book: Vehicles \u20b912,964 cr + CD \u20b96,065 cr = \u20b919,029 cr",
    "   \u2022 Target-book NPA: \u20b9404 cr @ 2.12% \u00b7 Annual write-offs: \u20b9437 cr",
    "   \u2022 Dealer advances: 0.23% NPA \u2014 dealers repay (fairness anchor)",
    "",
    "Conservative cash trajectory (front-loads full enterprise build in Y1):",
    "   \u2022 Phase 1 Pilot (0\u20133 mo): Shadow run on existing cloud, zero capex",
    "   \u2022 Y1: -\u20b912 cr (worst-case: \u20b99.5 cr full IT build/SOP absorbed in Y1 cash)",
    "   \u2022 Y2: +\u20b94 cr (breakeven reached as 52% of dealer book activates)",
    "   \u2022 Y3: +\u20b918 cr/yr net recurring profit (\u20b930 cr fraud saved \u2212 \u20b912 cr opex)",
    "   \u2022 Y4-Y5: +\u20b928 cr/yr (extended to \u20b928,111 cr book incl. tractors & PL)",
    "",
    "5-Yr Base NPV \u20b9+39 cr \u00b7 Upside \u20b9+64 cr \u00b7 89% P(NPV>0) \u00b7 4.2\u00d7 ROI",
], w=6.4, size=11)
add_image(s, "g5_financial.png", left=7.0)
footer(s, "[PUBLIC + A#] PROPOSAL/financial_model.py \u00b7 audited Note-18 anchors \u00b7 every assumption labelled")

# ---- S12 fairness -----------------------------------------------------------
s = add_slide("Rural fairness is a feature, not a caveat", "[PUBLIC + COLAB]")
bullets(s, [
    "__BIG__Creditworthiness ≠ fraud risk.",
    "A thin-file honest borrower is not a fraud suspect.",
    "",
    "Tested three independent ways, same verdict:",
    "   1. TVS’s audited books: dealer advances 0.23% NPA (FY25: 0.37%) — dealers repay.",
    "   2. Our ablation: dealer link carries the LEAST detection signal (−0.10 AUCPR).",
    "   3. Entity-resolution audit: 34 false merges found → why we never auto-reject on a link alone.",
    "",
    "__BIG__Three tests, one answer: the dealer is the conduit, not the cause.",
    "",
    "DPDP 2023: hashed identifiers, no PII in graph edges, purpose-limited. Explainable by construction.",
], w=6.4, size=13)
add_image(s, "g8_dealer_fairness.png", left=7.0)
footer(s, "[PUBLIC] audited AR Note 18 · [COLAB] ablation · [REPO] er_in_path.py · DPDP 2023")

# ---- S13 roadmap ------------------------------------------------------------
s = add_slide("Roadmap & change management", "[RICE]")
bullets(s, [
    "__BIG__Phase 1 (0–3 mo) · Shadow Rollout:",
    "One TN two-wheeler region. Backtest lead time on TVS’s own historical rings. NO automated blocking.",
    "",
    "Phase 2 (3–6 mo) · Advisory → Active:",
    "Link-freezes in underwriting. Analyst console. Dealer-fair scores.",
    "",
    "Phase 3 (6–12 mo) · Scale across India:",
    "Consumer durable, used car, personal loan, tractor. Adversarial red-team. Fairness monitoring.",
    "",
    "__BIG__Change management:",
    "2-week risk-analyst training · dealer SOP · kill-switch + rule fallback · pre-registered success metrics.",
], size=14)
footer(s, "Shadow-first rollout per lending-operations practice")

# ---- S14 reproducibility ----------------------------------------------------
s = add_slide("Run it yourself — and hold us to it", "[REPO]")
bullets(s, [
    "__BIG__30 seconds to reproduce:",
    "pip install -r requirements.txt",
    "python scripts/run_v1.py --profile SMOKE",
    "",
    "Colab notebook with all executed numbers (0.754, GNNs, YelpChi, ablations, leakage audits).",
    "",
    "No public dataset has ring-level lending fraud ground truth. Our synthetic pipeline is the only honest way to prove the methodology before TVS validates it on real rings in Phase 1.",
    "",
    "Three verified citations (AAAI-21 loan graphs, ICAIF-20 subgraphs, MDPI-25 causal hypergraphs).",
    "",
    "__BIG__Round 3 promise: in the finale, we run the engine live. Every chart on these slides will regenerate in front of you in 30 seconds.",
], size=14)
footer(s, "JALE repo  ·  colab/JALE_Colab_Training_ran.ipynb  ·  reports/*.json")

prs.save(OUT)
print(f"deck written: {OUT}")